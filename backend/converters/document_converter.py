import os
import csv
from pathlib import Path
from xml.sax.saxutils import escape as xml_escape


class DocumentConverter:
    SUPPORTED_INPUTS = {'.docx', '.xlsx', '.pptx', '.txt', '.csv'}

    def __init__(self, file_path):
        self.file_path = file_path
        self.filename_stem = Path(file_path).stem
        self.file_ext = Path(file_path).suffix.lower()
        self.output_dir = self._get_output_dir()

    def _get_output_dir(self):
        output_dir = os.path.join(os.path.dirname(self.file_path), '..', 'output')
        output_dir = os.path.abspath(output_dir)
        os.makedirs(output_dir, exist_ok=True)
        return output_dir

    def _build_pdf(self, story, output_path, landscape=False):
        """Constrói PDF usando reportlab com a lista de elementos fornecida."""
        from reportlab.lib.pagesizes import A4, landscape as rl_landscape
        from reportlab.platypus import SimpleDocTemplate
        from reportlab.lib.units import cm

        pagesize = rl_landscape(A4) if landscape else A4
        doc = SimpleDocTemplate(
            output_path, pagesize=pagesize,
            rightMargin=2 * cm, leftMargin=2 * cm,
            topMargin=2 * cm, bottomMargin=2 * cm
        )
        doc.build(story)

    # ─────────────── DOCX ───────────────

    def docx_to_txt(self):
        """Converte DOCX para TXT."""
        try:
            from docx import Document
            doc = Document(self.file_path)
            text = '\n'.join(p.text for p in doc.paragraphs)
            output_path = os.path.join(self.output_dir, f"{self.filename_stem}.txt")
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(text)
            return [output_path], None
        except Exception as e:
            return None, f"Erro ao converter DOCX para TXT: {str(e)}"

    def docx_to_pdf(self):
        """Converte DOCX para PDF (extração de texto via reportlab)."""
        try:
            from docx import Document
            from reportlab.lib.styles import getSampleStyleSheet
            from reportlab.platypus import Paragraph, Spacer
            from reportlab.lib.units import cm

            doc = Document(self.file_path)
            styles = getSampleStyleSheet()
            story = []

            for para in doc.paragraphs:
                if para.text.strip():
                    style_name = 'Heading1' if para.style.name.startswith('Heading') else 'Normal'
                    story.append(Paragraph(xml_escape(para.text), styles[style_name]))
                    story.append(Spacer(1, 0.3 * cm))

            output_path = os.path.join(self.output_dir, f"{self.filename_stem}.pdf")
            self._build_pdf(story, output_path)
            return [output_path], None
        except Exception as e:
            return None, f"Erro ao converter DOCX para PDF: {str(e)}"

    # ─────────────── XLSX ───────────────

    def xlsx_to_csv(self):
        """Converte XLSX para CSV (uma planilha = um arquivo CSV)."""
        try:
            from openpyxl import load_workbook
            wb = load_workbook(self.file_path, read_only=True, data_only=True)
            output_files = []

            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                safe = "".join(c for c in sheet_name if c.isalnum() or c in (' ', '-', '_')).strip()
                output_path = os.path.join(self.output_dir, f"{self.filename_stem}_{safe}.csv")
                with open(output_path, 'w', newline='', encoding='utf-8') as f:
                    writer = csv.writer(f)
                    for row in ws.iter_rows(values_only=True):
                        writer.writerow(['' if v is None else v for v in row])
                output_files.append(output_path)

            wb.close()
            return output_files, None
        except Exception as e:
            return None, f"Erro ao converter XLSX para CSV: {str(e)}"

    def xlsx_to_txt(self):
        """Converte XLSX para TXT."""
        try:
            from openpyxl import load_workbook
            wb = load_workbook(self.file_path, read_only=True, data_only=True)
            output_path = os.path.join(self.output_dir, f"{self.filename_stem}.txt")
            with open(output_path, 'w', encoding='utf-8') as f:
                for sheet_name in wb.sheetnames:
                    ws = wb[sheet_name]
                    f.write(f"=== {sheet_name} ===\n")
                    for row in ws.iter_rows(values_only=True):
                        f.write('\t'.join('' if v is None else str(v) for v in row) + '\n')
                    f.write('\n')
            wb.close()
            return [output_path], None
        except Exception as e:
            return None, f"Erro ao converter XLSX para TXT: {str(e)}"

    def xlsx_to_pdf(self):
        """Converte XLSX para PDF com tabela."""
        try:
            from openpyxl import load_workbook
            from reportlab.lib.styles import getSampleStyleSheet
            from reportlab.platypus import Table, TableStyle, Spacer, Paragraph
            from reportlab.lib import colors
            from reportlab.lib.units import cm

            wb = load_workbook(self.file_path, read_only=True, data_only=True)
            styles = getSampleStyleSheet()
            story = []

            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                story.append(Paragraph(xml_escape(f"Planilha: {sheet_name}"), styles['Heading2']))
                story.append(Spacer(1, 0.3 * cm))

                data = []
                for row in ws.iter_rows(values_only=True):
                    data.append(['' if v is None else str(v) for v in row])

                if data:
                    t = Table(data, repeatRows=1)
                    t.setStyle(TableStyle([
                        ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#334155')),
                        ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                        ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                        ('FONTSIZE', (0, 0), (-1, -1), 8),
                        ('GRID', (0, 0), (-1, -1), 0.5, colors.HexColor('#e2e8f0')),
                        ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.HexColor('#f8fafc')]),
                    ]))
                    story.append(t)
                story.append(Spacer(1, 0.5 * cm))

            wb.close()
            output_path = os.path.join(self.output_dir, f"{self.filename_stem}.pdf")
            self._build_pdf(story, output_path, landscape=True)
            return [output_path], None
        except Exception as e:
            return None, f"Erro ao converter XLSX para PDF: {str(e)}"

    # ─────────────── PPTX ───────────────

    def pptx_to_txt(self):
        """Extrai texto de todos os slides PPTX."""
        try:
            from pptx import Presentation
            prs = Presentation(self.file_path)
            output_path = os.path.join(self.output_dir, f"{self.filename_stem}.txt")
            with open(output_path, 'w', encoding='utf-8') as f:
                for i, slide in enumerate(prs.slides, 1):
                    f.write(f"=== Slide {i} ===\n")
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            f.write(shape.text + '\n')
                    f.write('\n')
            return [output_path], None
        except Exception as e:
            return None, f"Erro ao converter PPTX para TXT: {str(e)}"

    def pptx_to_pdf(self):
        """Converte PPTX para PDF (texto dos slides via reportlab)."""
        try:
            from pptx import Presentation
            from reportlab.lib.styles import getSampleStyleSheet
            from reportlab.platypus import Paragraph, Spacer
            from reportlab.lib.units import cm

            prs = Presentation(self.file_path)
            styles = getSampleStyleSheet()
            story = []

            for i, slide in enumerate(prs.slides, 1):
                story.append(Paragraph(f"Slide {i}", styles['Heading1']))
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        story.append(Paragraph(xml_escape(shape.text), styles['Normal']))
                        story.append(Spacer(1, 0.2 * cm))
                story.append(Spacer(1, 0.6 * cm))

            output_path = os.path.join(self.output_dir, f"{self.filename_stem}.pdf")
            self._build_pdf(story, output_path, landscape=True)
            return [output_path], None
        except Exception as e:
            return None, f"Erro ao converter PPTX para PDF: {str(e)}"

    # ─────────────── TXT ───────────────

    def txt_to_pdf(self):
        """Converte TXT para PDF."""
        try:
            from reportlab.lib.styles import getSampleStyleSheet
            from reportlab.platypus import Paragraph, Spacer
            from reportlab.lib.units import cm

            with open(self.file_path, 'r', encoding='utf-8', errors='replace') as f:
                lines = f.readlines()

            styles = getSampleStyleSheet()
            story = []
            for line in lines:
                text = line.rstrip('\n')
                story.append(Paragraph(xml_escape(text) if text else '&nbsp;', styles['Normal']))

            output_path = os.path.join(self.output_dir, f"{self.filename_stem}.pdf")
            self._build_pdf(story, output_path)
            return [output_path], None
        except Exception as e:
            return None, f"Erro ao converter TXT para PDF: {str(e)}"

    def txt_to_docx(self):
        """Converte TXT para DOCX."""
        try:
            from docx import Document
            with open(self.file_path, 'r', encoding='utf-8', errors='replace') as f:
                lines = f.readlines()
            doc = Document()
            for line in lines:
                doc.add_paragraph(line.rstrip('\n'))
            output_path = os.path.join(self.output_dir, f"{self.filename_stem}.docx")
            doc.save(output_path)
            return [output_path], None
        except Exception as e:
            return None, f"Erro ao converter TXT para DOCX: {str(e)}"

    # ─────────────── CSV ───────────────

    def csv_to_xlsx(self):
        """Converte CSV para XLSX."""
        try:
            from openpyxl import Workbook
            wb = Workbook()
            ws = wb.active
            ws.title = "Dados"
            with open(self.file_path, 'r', encoding='utf-8', errors='replace', newline='') as f:
                reader = csv.reader(f)
                for row in reader:
                    ws.append(row)
            output_path = os.path.join(self.output_dir, f"{self.filename_stem}.xlsx")
            wb.save(output_path)
            return [output_path], None
        except Exception as e:
            return None, f"Erro ao converter CSV para XLSX: {str(e)}"

    def csv_to_txt(self):
        """Copia CSV como TXT (preservando conteúdo)."""
        try:
            with open(self.file_path, 'r', encoding='utf-8', errors='replace') as f:
                content = f.read()
            output_path = os.path.join(self.output_dir, f"{self.filename_stem}.txt")
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(content)
            return [output_path], None
        except Exception as e:
            return None, f"Erro ao converter CSV para TXT: {str(e)}"
