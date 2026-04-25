import os
import io
from pathlib import Path
from pdf2image import convert_from_path
from PyPDF2 import PdfReader
from PIL import Image
from docx import Document
from docx.shared import Pt
from openpyxl import Workbook
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.util import Inches, Pt as PptPt

class PDFConverter:
    def __init__(self, pdf_path):
        self.pdf_path = pdf_path
        self.pdf_filename = Path(pdf_path).stem
    
    def to_images(self, format='PNG', dpi=150):
        """Converte PDF para imagens (PNG ou JPG)."""
        try:
            images = convert_from_path(self.pdf_path, dpi=dpi)
            output_files = []
            
            for i, image in enumerate(images):
                output_filename = f"{self.pdf_filename}_page_{i+1:03d}.{format.lower()}"
                output_path = os.path.join('./temp/output', output_filename)
                
                if format.upper() == 'JPG':
                    # Converter RGBA para RGB para JPG
                    if image.mode in ('RGBA', 'LA', 'P'):
                        rgb_image = Image.new('RGB', image.size, (255, 255, 255))
                        rgb_image.paste(image, mask=image.split()[-1] if image.mode == 'RGBA' else None)
                        rgb_image.save(output_path, 'JPEG', quality=95)
                    else:
                        image.save(output_path, 'JPEG', quality=95)
                else:
                    image.save(output_path, 'PNG')
                
                output_files.append(output_path)
            
            return output_files, None
        except Exception as e:
            return None, f"Erro ao converter para imagens: {str(e)}"
    
    def to_text(self):
        """Converte PDF para TXT."""
        try:
            text = ""
            with open(self.pdf_path, 'rb') as file:
                pdf_reader = PdfReader(file)
                for page_num, page in enumerate(pdf_reader.pages, 1):
                    text += f"--- Página {page_num} ---\n"
                    text += page.extract_text()
                    text += "\n\n"
            
            output_filename = f"{self.pdf_filename}.txt"
            output_path = os.path.join('./temp/output', output_filename)
            
            with open(output_path, 'w', encoding='utf-8') as f:
                f.write(text)
            
            return [output_path], None
        except Exception as e:
            return None, f"Erro ao converter para TXT: {str(e)}"
    
    def to_docx(self):
        """Converte PDF para DOCX."""
        try:
            doc = Document()
            doc.add_heading(self.pdf_filename, 0)
            
            with open(self.pdf_path, 'rb') as file:
                pdf_reader = PdfReader(file)
                for page_num, page in enumerate(pdf_reader.pages, 1):
                    doc.add_heading(f"Página {page_num}", level=1)
                    text = page.extract_text()
                    doc.add_paragraph(text)
            
            output_filename = f"{self.pdf_filename}.docx"
            output_path = os.path.join('./temp/output', output_filename)
            doc.save(output_path)
            
            return [output_path], None
        except Exception as e:
            return None, f"Erro ao converter para DOCX: {str(e)}"
    
    def to_xlsx(self):
        """Converte PDF para XLSX."""
        try:
            wb = Workbook()
            ws = wb.active
            ws.title = "PDF Data"
            
            with open(self.pdf_path, 'rb') as file:
                pdf_reader = PdfReader(file)
                row = 1
                
                for page_num, page in enumerate(pdf_reader.pages, 1):
                    text = page.extract_text()
                    lines = text.split('\n')
                    
                    ws[f'A{row}'] = f"Página {page_num}"
                    ws[f'A{row}'].font = ws[f'A{row}'].font.copy()
                    row += 1
                    
                    for line in lines:
                        if line.strip():
                            ws[f'A{row}'] = line
                            row += 1
                    
                    row += 1
            
            output_filename = f"{self.pdf_filename}.xlsx"
            output_path = os.path.join('./temp/output', output_filename)
            wb.save(output_path)
            
            return [output_path], None
        except Exception as e:
            return None, f"Erro ao converter para XLSX: {str(e)}"
    
    def to_pptx(self):
        """Converte PDF para PPTX (uma página por slide com imagem)."""
        try:
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
            
            images = convert_from_path(self.pdf_path, dpi=150)
            
            for i, image in enumerate(images):
                # Salvar imagem temporária
                temp_img_path = f'./temp/temp_page_{i}.png'
                image.save(temp_img_path, 'PNG')
                
                # Adicionar slide
                blank_slide_layout = prs.slide_layouts[6]
                slide = prs.slides.add_slide(blank_slide_layout)
                
                # Adicionar imagem ao slide
                left = Inches(0)
                top = Inches(0)
                pic = slide.shapes.add_picture(temp_img_path, left, top, width=prs.slide_width, height=prs.slide_height)
                
                # Limpar imagem temporária
                os.remove(temp_img_path)
            
            output_filename = f"{self.pdf_filename}.pptx"
            output_path = os.path.join('./temp/output', output_filename)
            prs.save(output_path)
            
            return [output_path], None
        except Exception as e:
            return None, f"Erro ao converter para PPTX: {str(e)}"
